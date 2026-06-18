from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import RGBColor
from app.utils.geometry import px_to_inches


def _rgb(hex_color):
    h=(hex_color or '#000000').lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))


def _add_picture(slide, job_root: Path, asset_path: str, x, y, w, h) -> None:
    path = job_root / asset_path
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def build_pptx(manifest, job_root: Path, output_path: Path) -> Path:
    prs=Presentation(); prs.slide_width=Inches(manifest.slide.width_in); prs.slide_height=Inches(manifest.slide.height_in)
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    for e in manifest.elements:
        x,y,w,h = px_to_inches(e.bbox_px, manifest.source.width_px, manifest.source.height_px, manifest.slide.width_in, manifest.slide.height_in)
        if e.type == 'background' and e.asset_path:
            _add_picture(slide, job_root, e.asset_path, x, y, w, h)
        elif e.type == 'text':
            tb=slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
            tf = tb.text_frame
            tf.margin_left = Pt(e.style.get('margin_left', 0))
            tf.margin_right = Pt(e.style.get('margin_right', 0))
            tf.margin_top = Pt(e.style.get('margin_top', 0))
            tf.margin_bottom = Pt(e.style.get('margin_bottom', 0))
            tf.word_wrap = False
            tf.vertical_anchor = MSO_ANCHOR.TOP
            p=tf.paragraphs[0]
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            run=p.add_run(); run.text=e.text or ''; run.font.size=Pt(e.style.get('font_size',14)); run.font.name=e.style.get('font_family','Arial'); run.font.color.rgb=_rgb(e.style.get('color','#1F2937')); run.font.bold=bool(e.style.get('bold',False))
        elif e.type == 'shape':
            shape_type=MSO_SHAPE.ROUNDED_RECTANGLE if e.shape=='rounded_rect' else MSO_SHAPE.RECTANGLE
            shp=slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h)); shp.fill.solid(); shp.fill.fore_color.rgb=_rgb(e.style.get('fill','#FFFFFF')); shp.line.color.rgb=_rgb(e.style.get('stroke','#D1D5DB'))
        elif e.type in ('image', 'icon', 'chart', 'table') and e.asset_path:
            _add_picture(slide, job_root, e.asset_path, x, y, w, h)
        elif e.type in ('line','arrow'):
            slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x+w), Inches(y+h))
    prs.save(output_path); return output_path
